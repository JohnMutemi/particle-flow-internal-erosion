from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Advanced CFD-DEM Coupling Framework"
    subtitle.text = "Project Presentation\nGeotechnical Applications"

def create_overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Project Overview"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Key Components:"
    
    p = tf.add_paragraph()
    p.text = "• Validation Manager"
    p = tf.add_paragraph()
    p.text = "• Statistical Analysis"
    p = tf.add_paragraph()
    p.text = "• Sensitivity Analysis"
    p = tf.add_paragraph()
    p.text = "• Case Studies"

def create_visualization_slide(prs, title, image_path, left, top):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    title_shape.text = title
    
    slide.shapes.add_picture(image_path, left, top, width=Inches(6))

def main():
    prs = Presentation()
    
    # Title slide
    create_title_slide(prs)
    
    # Overview slide
    create_overview_slide(prs)
    
    # Framework Components
    create_visualization_slide(prs, "Fluid-Particle Coupling", 
                             "results/coupling_forces.png",
                             Inches(1), Inches(2))
    
    create_visualization_slide(prs, "Coarse-Grained Approach",
                             "results/coarse_grained.png",
                             Inches(1), Inches(2))
    
    # Statistical Analysis
    create_visualization_slide(prs, "Fluid Flow Statistics",
                             "results/validation/statistics_fluid.png",
                             Inches(1), Inches(2))
    
    create_visualization_slide(prs, "Particle Behavior Analysis",
                             "results/validation/statistics_particles.png",
                             Inches(1), Inches(2))
    
    # Sensitivity Analysis
    create_visualization_slide(prs, "Sensitivity Analysis Results",
                             "results/validation/sensitivity_erosion_rate_coefficient.png",
                             Inches(1), Inches(2))
    
    # Validation Results
    create_visualization_slide(prs, "Validation Results",
                             "results/validation/comparison_velocity.png",
                             Inches(1), Inches(2))
    
    # Case Studies
    create_visualization_slide(prs, "Bond Degradation Model",
                             "results/bond_degradation.png",
                             Inches(1), Inches(2))
    
    # Save the presentation
    prs.save('CFD_DEM_Presentation.pptx')

if __name__ == '__main__':
    main() 